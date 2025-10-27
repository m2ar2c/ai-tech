#!/usr/bin/env python3
"""Convert doubao-01.html into a structured PowerPoint deck."""
from __future__ import annotations

import re
from pathlib import Path
from typing import Iterable, List

from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.util import Pt

BASE_DIR = Path(__file__).resolve().parent
INPUT_HTML = BASE_DIR / "doubao-01.html"
OUTPUT_PPTX = BASE_DIR / "doubao-01.pptx"

FONT_NAME = "Microsoft YaHei"
TITLE_FONT_SIZE = Pt(42)
COVER_TITLE_FONT_SIZE = Pt(48)
SUBTITLE_FONT_SIZE = Pt(28)
BODY_FONT_SIZE = Pt(24)
DEFAULT_CHUNK_SIZE = 8

SKIP_TAGS = {"script", "style", "canvas", "svg", "button"}


def clean_text(text: str) -> str:
    """Normalize whitespace in extracted text."""
    return re.sub(r"\s+", " ", text).strip()


def is_major_h3(tag: Tag | None) -> bool:
    """Identify subsection headings styled as major titles (text-xl)."""
    if not isinstance(tag, Tag) or tag.name != "h3":
        return False
    classes = tag.get("class") or []
    return "text-xl" in classes


def contains_major_h3(tag: Tag | None) -> bool:
    if not isinstance(tag, Tag):
        return False
    return tag.find(is_major_h3) is not None


def extract_texts(node: Tag | NavigableString | None) -> List[str]:
    """Recursively extract meaningful text fragments from a node."""
    texts: List[str] = []
    if node is None:
        return texts

    if isinstance(node, NavigableString):
        text = clean_text(str(node))
        if text:
            texts.append(text)
        return texts

    if not isinstance(node, Tag):
        return texts

    if node.name in SKIP_TAGS:
        return texts

    if node.name == "li":
        text = clean_text(node.get_text(" ", strip=True))
        if text:
            texts.append(text)
        return texts

    if node.name in {"p", "blockquote", "h4", "h5", "h6", "dt", "dd"}:
        text = clean_text(node.get_text(" ", strip=True))
        if text:
            texts.append(text)
        return texts

    if node.name in {"ul", "ol"}:
        for li in node.find_all("li", recursive=False):
            texts.extend(extract_texts(li))
        return texts

    if node.name == "table":
        for tr in node.find_all("tr"):
            cells = [clean_text(td.get_text(" ", strip=True)) for td in tr.find_all(["th", "td"])]
            filtered = [cell for cell in cells if cell]
            if filtered:
                texts.append(" | ".join(filtered))
        return texts

    if node.name == "dl":
        for child in node.find_all(["dt", "dd"], recursive=False):
            texts.extend(extract_texts(child))
        return texts

    # Skip major subsection headings when traversing containers
    if is_major_h3(node):
        return texts

    for child in node.children:
        texts.extend(extract_texts(child))
    return texts


def extract_texts_before_major(tag: Tag) -> List[str]:
    """Extract text from a tag before the first major h3 descendant."""
    fragments: List[str] = []
    for child in tag.children:
        if isinstance(child, NavigableString):
            text = clean_text(str(child))
            if text:
                fragments.append(text)
            continue

        if not isinstance(child, Tag):
            continue

        if is_major_h3(child):
            break

        if contains_major_h3(child):
            fragments.extend(extract_texts_before_major(child))
            break

        fragments.extend(extract_texts(child))
    return fragments


def post_process_bullets(items: Iterable[str]) -> List[str]:
    """Clean, deduplicate, and combine related bullet entries."""
    cleaned = [clean_text(item) for item in items if clean_text(item)]
    combined: List[str] = []
    idx = 0
    while idx < len(cleaned):
        current = cleaned[idx]
        if idx + 1 < len(cleaned):
            nxt = cleaned[idx + 1]
            if len(current) <= 12 and not re.search(r"[。！？!?]$", current) and len(nxt) > 15:
                combined.append(f"{current}：{nxt}")
                idx += 2
                continue
        combined.append(current)
        idx += 1

    deduped: List[str] = []
    seen = set()
    for item in combined:
        if item not in seen:
            deduped.append(item)
            seen.add(item)
    return deduped


def get_content_blocks(h2: Tag) -> List[Tag]:
    """Return the block-level elements that contain actual section content."""
    inner_wrapper = h2.find_parent(lambda t: isinstance(t, Tag) and t.name == "div" and t.get("class") and "max-w-4xl" in t.get("class"))
    if not inner_wrapper:
        return []

    header_container = h2.parent
    blocks: List[Tag] = []
    capture = False
    for child in inner_wrapper.children:
        if child is header_container:
            capture = True
            continue
        if not capture:
            continue
        if isinstance(child, Tag):
            blocks.append(child)
    return blocks


def gather_intro_paragraphs(h2: Tag) -> List[str]:
    blocks = get_content_blocks(h2)
    intro_texts: List[str] = []
    for block in blocks:
        if contains_major_h3(block):
            intro_texts.extend(extract_texts_before_major(block))
            break
        intro_texts.extend(extract_texts(block))
    return post_process_bullets(intro_texts)


def gather_subsection_content(h3_tag: Tag) -> List[str]:
    fragments: List[str] = []
    for sibling in h3_tag.next_siblings:
        if isinstance(sibling, Tag):
            if is_major_h3(sibling):
                break
            fragments.extend(extract_texts(sibling))
        elif isinstance(sibling, NavigableString):
            text = clean_text(str(sibling))
            if text:
                fragments.append(text)
    return post_process_bullets(fragments)


def chunk_list(items: List[str], size: int) -> Iterable[List[str]]:
    for idx in range(0, len(items), size):
        yield items[idx : idx + size]


def add_single_bullet_slide(prs: Presentation, title: str, bullet_items: List[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]

    title_shape.text = title
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.font.name = FONT_NAME
    title_para.font.size = TITLE_FONT_SIZE

    text_frame = body_shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True

    for idx, text in enumerate(bullet_items):
        paragraph = text_frame.add_paragraph() if idx else text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.level = 0
        paragraph.font.name = FONT_NAME
        paragraph.font.size = BODY_FONT_SIZE


def add_bullet_slides(prs: Presentation, title: str, items: List[str], chunk_size: int = DEFAULT_CHUNK_SIZE) -> None:
    valid_items = [item for item in items if item]
    if not valid_items:
        return
    total_chunks = (len(valid_items) + chunk_size - 1) // chunk_size
    for index, chunk in enumerate(chunk_list(valid_items, chunk_size), start=1):
        slide_title = title if total_chunks == 1 else f"{title}（{index}/{total_chunks}）"
        add_single_bullet_slide(prs, slide_title, chunk)


def build_presentation(soup: BeautifulSoup) -> Presentation:
    prs = Presentation()

    # Cover slide
    cover_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = cover_slide.shapes.title
    subtitle_shape = cover_slide.placeholders[1]

    hero_section = soup.select_one("section.bg-gradient-tech")
    hero_title_tag = hero_section.find("h1") if hero_section else soup.find("h1")
    hero_subtitle_tag = None
    if hero_section:
        hero_subtitle_tag = hero_section.find("p")
    elif hero_title_tag:
        hero_subtitle_tag = hero_title_tag.find_next("p")

    cover_title_text = clean_text(hero_title_tag.get_text(" ", strip=True)) if hero_title_tag else "AI技术核心解析报告"
    title_shape.text = cover_title_text
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.font.name = FONT_NAME
    title_para.font.size = COVER_TITLE_FONT_SIZE

    subtitle_text = clean_text(hero_subtitle_tag.get_text(" ", strip=True)) if hero_subtitle_tag else "核心内容概览"
    subtitle_shape.text = subtitle_text
    subtitle_para = subtitle_shape.text_frame.paragraphs[0]
    subtitle_para.font.name = FONT_NAME
    subtitle_para.font.size = SUBTITLE_FONT_SIZE

    main_el = soup.find("main")
    if not main_el:
        return prs

    sections = [sec for sec in main_el.find_all("section", recursive=False)]

    # Agenda slide
    section_titles = []
    for sec in sections:
        h2 = sec.find("h2")
        if not h2:
            continue
        title = clean_text(h2.get_text(" ", strip=True))
        if title:
            section_titles.append(title)
    add_bullet_slides(prs, "目录", section_titles, chunk_size=9)

    for sec in sections:
        h2 = sec.find("h2")
        if not h2:
            continue
        sec_title = clean_text(h2.get_text(" ", strip=True))

        intro_items = gather_intro_paragraphs(h2)
        add_bullet_slides(prs, sec_title, intro_items)

        major_headings = sec.find_all(is_major_h3)
        seen_ids: set[int] = set()
        for h3_tag in major_headings:
            if id(h3_tag) in seen_ids:
                continue
            seen_ids.add(id(h3_tag))
            subsection_title = clean_text(h3_tag.get_text(" ", strip=True))
            content_items = gather_subsection_content(h3_tag)
            if not content_items:
                continue
            slide_title = f"{sec_title}｜{subsection_title}"
            add_bullet_slides(prs, slide_title, content_items)

    return prs


def main() -> None:
    if not INPUT_HTML.exists():
        raise FileNotFoundError(f"输入文件不存在: {INPUT_HTML}")

    html = INPUT_HTML.read_text(encoding="utf-8")
    soup = BeautifulSoup(html, "lxml")

    presentation = build_presentation(soup)
    presentation.save(OUTPUT_PPTX)
    print(f"已生成演示文稿: {OUTPUT_PPTX}")


if __name__ == "__main__":
    main()
